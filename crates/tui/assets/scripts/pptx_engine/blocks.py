"""pptx_engine/blocks.py — block renderers and drawing primitives.

Phase 1: table, richtext (runs), image (PNG/JPG with fit), chart dispatch,
drawing primitives (bg, rect, textbox, paragraph), dimension parsing.
"""

import os
import re
import sys

try:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    pass

from .charts import add_chart
from .mpl import render_mpl_chart
from .theme import _hex_to_rgb

# ── Dimension parsing ────────────────────────────────────────────────────

_DIM_RE = re.compile(r"^([\d.]+)\s*(in|pt|emu)$", re.IGNORECASE)


def parse_dim(s):
    """Parse a dimension string to EMU.

    Supported units: in (inches), pt (points), emu.
    Returns an EMU int.
    """
    s = str(s).strip()
    m = _DIM_RE.match(s)
    if not m:
        raise ValueError(f"invalid dimension: {s!r}")
    val = float(m.group(1))
    unit = m.group(2).lower()
    if unit == "in":
        return int(val * 914400)
    elif unit == "pt":
        return int(val * 12700)
    else:  # emu
        return int(val)


def _emu_to_inches(emu_val):
    """Convert EMU to float inches."""
    return emu_val / 914400.0


# ── Drawing primitives ───────────────────────────────────────────────────

def add_bg(slide, clr):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = clr


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def add_para(tf, text, size, color, font, bold=False, italic=False, align=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.space_after = Pt(int(size * 0.35))
    p.space_before = Pt(int(size * 0.1))
    if align is not None:
        p.alignment = align
    return p


# ── Table block ──────────────────────────────────────────────────────────

def add_table(slide, table_data, t, left=None, top=None, width=None):
    """Add a styled native table to the slide.

    Args:
        slide: pptx slide object.
        table_data: dict with 'headers' and 'rows'.
        t: theme dict with table colors.
        left, top, width: optional override (EMU). Defaults to full-width centered.

    Returns:
        Total height in EMU, or 0 if no data.
    """
    td = table_data
    headers = td.get("headers", [])
    rows = td.get("rows", [])
    if not headers and not rows:
        return 0

    n_cols = len(headers) or (len(rows[0]) if rows else 1)
    n_body_rows = len(rows)
    n_total_rows = n_body_rows + (1 if headers else 0)

    max_rows = 18
    if n_total_rows > max_rows:
        n_body_rows = max_rows - (1 if headers else 0)
        rows = rows[:n_body_rows]
        n_total_rows = max_rows

    row_h = Pt(26)
    total_h = row_h * n_total_rows + Pt(2)

    if left is None:
        left = Inches(0.6)
    if top is None:
        top = Inches(1.85)
    if width is None:
        width = Inches(12.133)

    tbl_shape = slide.shapes.add_table(n_total_rows, n_cols, left, top, width, total_h)
    tbl = tbl_shape.table

    col_w = int(width / n_cols)
    for ci in range(n_cols):
        tbl.columns[ci].width = col_w

    def _cell(ri, ci, text, color, bg_color, font_name, bold=False, size=11):
        cell = tbl.cell(ri, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = PP_ALIGN.LEFT
        cell.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

    ri = 0
    if headers:
        for ci, h in enumerate(headers):
            _cell(ri, ci, str(h), t["table_header_fg"], t["table_header_bg"],
                  t["font"], bold=True, size=12)
        ri += 1

    for i, row in enumerate(rows):
        bg = t["table_alt_bg"] if i % 2 else t["table_row_bg"]
        for ci in range(n_cols):
            val = str(row[ci]) if ci < len(row) else ""
            _cell(ri, ci, val, t["body"], bg, t["font"], size=11)
        ri += 1

    # Merge cells if specified: [[r1, c1, r2, c2], ...] (0-based, r1≤r2, c1≤c2)
    merges = td.get("merges", [])
    for m in merges:
        if len(m) == 4:
            r1, c1, r2, c2 = m
            if 0 <= r1 <= r2 < n_total_rows and 0 <= c1 <= c2 < n_cols:
                tbl.cell(r1, c1).merge(tbl.cell(r2, c2))

    return total_h


# ── Block dispatcher ─────────────────────────────────────────────────────

def render_block(slide, block, left, top, width, height, t):
    """Render a single block into its region.

    Supported block types:
        chart, table, richtext, image, text (fallback).
    """
    btype = block.get("type", "text")

    if btype == "chart":
        chart_data = block.get("chart")
        if chart_data:
            add_chart(slide, chart_data, left, top, width, height, t=t)

    elif btype == "table":
        table_data = block.get("table", block)
        add_table(slide, table_data, t, left=left, top=top, width=width)

    elif btype == "richtext":
        runs = block.get("runs", [])
        if runs:
            tf = add_textbox(slide, left, top, width, height)
            for i, run in enumerate(runs):
                text = str(run.get("t", ""))
                size = run.get("size", 14)
                color_hex = run.get("color")
                color = _hex_to_rgb(color_hex) if color_hex else t["body"]
                bold = run.get("bold", False)
                italic = run.get("italic", False)
                add_para(tf, text, size, color, t["font"],
                         bold=bold, italic=italic, first=(i == 0))

    elif btype == "image":
        path = block.get("path", "")
        if path:
            render_image(slide, block, left, top, width, height, t)

    elif btype == "mpl":
        mpl_data = block.get("mpl", block)
        png_path = render_mpl_chart(mpl_data)
        if png_path:
            img_block = {"path": png_path, "fit": block.get("fit", "contain")}
            render_image(slide, img_block, left, top, width, height, t)

    else:
        # Fallback: plain text
        text = block.get("text", "")
        if text:
            tf = add_textbox(slide, left, top, width, height)
            add_para(tf, str(text), 14, t["body"], t["font"], first=True)


# ── Image block ──────────────────────────────────────────────────────────

def _get_image_native_size(path):
    """Return (width_px, height_px) or None if PIL is unavailable."""
    try:
        from PIL import Image as PILImage
        with PILImage.open(path) as im:
            return im.size
    except Exception:
        return None


def render_image(slide, block, left, top, width, height, t):
    """Render an image block with fit=contain|cover support.

    Args:
        slide: pptx slide object.
        block: dict with 'path', optional 'fit' (contain|cover), 'max_height'.
        left, top, width, height: bounding box in EMU.
        t: theme dict.
    """
    path = block.get("path", "")
    if not path or not os.path.isfile(path):
        print(f"WARNING: image not found: {path!r}", file=sys.stderr)
        tf = add_textbox(slide, left, top, width, height)
        add_para(tf, f"[Image: {os.path.basename(path) or 'missing'}]",
                 12, t["muted"], t["font"], first=True)
        return

    native_size = _get_image_native_size(path)
    fit = block.get("fit", "contain")
    max_h_str = block.get("max_height")

    try:
        if native_size:
            img_w, img_h = native_size
            img_aspect = img_w / img_h if img_h > 0 else 1.0
            box_w = _emu_to_inches(width)
            box_h = _emu_to_inches(height)
            box_aspect = box_w / box_h if box_h > 0 else 1.0

            if fit == "contain":
                # Scale proportionally to fit within box, then center
                if img_aspect > box_aspect:
                    # Image wider — scale by width
                    pic_w = width
                    pic_h = int(width / img_aspect)
                else:
                    # Image taller — scale by height
                    pic_h = height
                    pic_w = int(height * img_aspect)
                # Center in box
                pic_left = left + (width - pic_w) // 2
                pic_top = top + (height - pic_h) // 2

            elif fit == "cover":
                # Scale proportionally to fill box (may overflow)
                if img_aspect > box_aspect:
                    # Image wider — scale by height to fill vertically
                    pic_h = height
                    pic_w = int(height * img_aspect)
                else:
                    # Image taller — scale by width to fill horizontally
                    pic_w = width
                    pic_h = int(width / img_aspect)
                # Center overflow
                pic_left = left + (width - pic_w) // 2
                pic_top = top + (height - pic_h) // 2

            else:
                # Unknown fit mode — stretch to fill
                pic_left, pic_top = left, top
                pic_w, pic_h = width, height

            pic = slide.shapes.add_picture(path, pic_left, pic_top, pic_w, pic_h)
        else:
            # No PIL — stretch to fill
            pic = slide.shapes.add_picture(path, left, top, width, height)

        # Apply max_height constraint
        if max_h_str:
            max_h_inches = _emu_to_inches(parse_dim(max_h_str))
            current_h = _emu_to_inches(pic.height)
            if current_h > max_h_inches:
                aspect = _emu_to_inches(pic.width) / current_h
                new_h = Inches(max_h_inches)
                new_w = int(new_h * aspect)
                pic.height = new_h
                pic.width = new_w

    except Exception as e:
        print(f"WARNING: failed to embed image {path!r}: {e}", file=sys.stderr)
        tf = add_textbox(slide, left, top, width, height)
        add_para(tf, f"[Image error: {e}]", 12, t["muted"], t["font"], first=True)
