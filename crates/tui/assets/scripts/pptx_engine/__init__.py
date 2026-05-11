"""pptx_engine — PPTX generation engine.

Entry point: build_presentation(payload) → Presentation

Phase 1 provides:
- Theme resolution with 4 presets + custom (theme.py)
- Grid layout engine: weight-based columns with gap/padding (layout.py)
- Block renderers: richtext (runs), image (PNG/JPG), table, chart (blocks.py + charts.py)
- Old pipeline compatibility: chart → table → bullets (layout.py)
- Combo chart (bar + line overlay, e.g. Pareto)
- §5.0 blocks/old-pipeline priority with error fallback
"""

import sys

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("ERROR: python-pptx not installed", file=sys.stderr)
    sys.exit(1)

from .theme import resolve_theme
from .layout import (
    SLIDE_W, SLIDE_H,
    build_cover, build_content, build_content_blocks, build_end,
)


def build_presentation(payload):
    """Build a complete Presentation from a JSON payload.

    Args:
        payload: dict with optional 'title', 'subtitle', 'theme',
                 and required 'slides' array.

    Returns:
        pptx.Presentation object ready for save().
    """
    # ── Theme ──
    try:
        t = resolve_theme(payload.get("theme"))
    except ValueError as e:
        print(f"ERROR: invalid theme — {e}", file=sys.stderr)
        sys.exit(1)

    doc_title = payload.get("title", "")
    doc_subtitle = payload.get("subtitle", "")
    slides_data = payload.get("slides", [])

    if not slides_data:
        print("ERROR: no slides data", file=sys.stderr)
        sys.exit(1)

    # ── Presentation setup ──
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    total = len(slides_data) + 1  # +1 for end slide
    if doc_title:
        total += 1  # +1 for cover

    # ── Cover ──
    if doc_title:
        build_cover(prs.slides.add_slide(prs.slide_layouts[6]),
                    doc_title, doc_subtitle, t)

    # ── Content slides ──
    base = 1 if doc_title else 0
    for i, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # §5.0: blocks pipeline takes priority over old pipeline
        blocks = sd.get("blocks")
        if blocks is not None and len(blocks) > 0:
            build_content_blocks(slide, sd, i + 1 + base, total, t)
        else:
            build_content(slide, sd, i + 1 + base, total, t)

    # ── End slide ──
    build_end(prs.slides.add_slide(prs.slide_layouts[6]),
              doc_title or "", t)

    return prs
