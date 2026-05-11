"""pptx_engine/template.py — template-driven PPTX generation.

Phase 2: load a base .pptx, fill text/image placeholders, clone slides.

Design constraints (§3.2):
  - Text: {{key}} in shape text → string replace
  - Image: shape name match → replace with image file
  - Slide clone: duplicate template slide, then fill
  - NOT doing: table cell filling in templates (python-pptx limitations)
"""

import os
import sys
import copy

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    pass


def load_template(template_path):
    """Load a .pptx file as template."""
    if not os.path.isfile(template_path):
        print(f"ERROR: template not found: {template_path!r}", file=sys.stderr)
        return None
    try:
        return Presentation(template_path)
    except Exception as e:
        print(f"ERROR: failed to load template: {e}", file=sys.stderr)
        return None


def fill_text_placeholders(slide, values):
    """Replace {{key}} patterns in all text shapes on a slide.

    Args:
        slide: pptx slide object.
        values: dict of {key: replacement_string}.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, val in values.items():
                        placeholder = f"{{{{{key}}}}}"
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, str(val))


def fill_image_placeholders(slide, image_map):
    """Replace image placeholder shapes with actual images.

    Args:
        slide: pptx slide object.
        image_map: dict of {shape_name: image_path}.
    """
    for shape in slide.shapes:
        if shape.name in image_map:
            img_path = image_map[shape.name]
            if os.path.isfile(img_path):
                try:
                    # Replace the placeholder shape with an image
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    # Remove old shape and add picture
                    sp = shape._element
                    sp.getparent().remove(sp)
                    slide.shapes.add_picture(img_path, left, top, width, height)
                except Exception as e:
                    print(f"WARNING: failed to replace image '{shape.name}': {e}",
                          file=sys.stderr)


def build_from_template(template_path, slide_fills, output_path=None):
    """Build a presentation from a template by filling placeholders.

    Args:
        template_path: path to base .pptx file.
        slide_fills: list of dicts, each with:
            - slide_index (int): 0-based index in template to clone
            - text: {key: value} for {{key}} placeholders
            - images: {shape_name: image_path} for image replacements
            - clone_count (int, optional): duplicate this slide N times
        output_path: optional save path.

    Returns:
        pptx.Presentation object, or None on failure.
    """
    prs = load_template(template_path)
    if prs is None:
        return None

    # Build new presentation from filled slides
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height

    for fill in slide_fills:
        slide_idx = fill.get("slide_index", 0)
        clone_count = fill.get("clone_count", 1)
        text_vals = fill.get("text", {})
        img_map = fill.get("images", {})

        if slide_idx >= len(prs.slides):
            print(f"WARNING: slide_index {slide_idx} out of range", file=sys.stderr)
            continue

        for _ in range(clone_count):
            # Clone template slide layout to new presentation
            template_slide = prs.slides[slide_idx]
            # Use blank layout and copy shapes manually
            new_slide = new_prs.slides.add_slide(new_prs.slide_layouts[6])

            # Copy all shapes from template slide
            _copy_shapes(template_slide, new_slide)

            # Fill placeholders
            fill_text_placeholders(new_slide, text_vals)
            fill_image_placeholders(new_slide, img_map)

    if output_path:
        new_prs.save(output_path)

    return new_prs


def _copy_shapes(source_slide, dest_slide):
    """Copy all shapes from source slide to destination slide.

    python-pptx doesn't have a native slide clone API, so we copy
    shape elements via XML manipulation.
    """
    import copy
    from lxml import etree

    # Copy slide background
    bg = source_slide.background
    if bg.fill.type is not None:
        dest_slide.background.fill.solid()
        try:
            dest_slide.background.fill.fore_color.rgb = bg.fill.fore_color.rgb
        except Exception:
            pass

    # Copy shapes
    for shape in source_slide.shapes:
        el = copy.deepcopy(shape._element)
        dest_slide.shapes._spTree.append(el)


def make_template_data(slide_fills):
    """Build a payload-compatible dict for use with build_presentation().

    This wraps template generation into the standard JSON pipeline,
    returning a payload where slides are generated from a template.

    Args:
        slide_fills: same as build_from_template's slide_fills.

    Returns:
        dict with 'template' key for use by build_presentation().
    """
    return {"template": True, "slide_fills": slide_fills}
