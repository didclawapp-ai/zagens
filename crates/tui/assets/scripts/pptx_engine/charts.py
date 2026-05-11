"""pptx_engine/charts.py — native OOXML chart helpers.

Migrated from write_pptx.py. Provides CHART_TYPES and _add_chart().
Phase 1 will add multi-chart-per-slide support and combo (bar+line) strategy.
"""

try:
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Pt
except ImportError:
    CategoryChartData = None
    XyChartData = None
    XL_CHART_TYPE = None

CHART_TYPES = {
    "bar":            XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line":           XL_CHART_TYPE.LINE_MARKERS,
    "pie":            XL_CHART_TYPE.PIE,
    "stacked_bar":    XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_bar_pct": XL_CHART_TYPE.COLUMN_STACKED_100,
    "area":           XL_CHART_TYPE.AREA,
    "scatter":        XL_CHART_TYPE.XY_SCATTER,
    "donut":          XL_CHART_TYPE.DOUGHNUT,
    "combo":          None,  # handled specially — bar + line overlay
}


def _try_category_axis(chart):
    """Return category/date axis or None (pie/donut/scatter layouts vary)."""
    try:
        return chart.category_axis
    except ValueError:
        return None


def _try_value_axis(chart):
    """Return primary value axis or None."""
    try:
        return chart.value_axis
    except ValueError:
        return None


def _apply_rgb_to_textframe(tf, rgb):
    """Force RGB on defRPr, every paragraph, and every run.

    Word inserts theme-based runs; paragraph-level color alone is often ignored
    when the slide uses a dark background.
    """
    if tf is None:
        return
    try:
        tf.paragraphs[0].font.color.rgb = rgb
    except (AttributeError, TypeError, ValueError):
        pass
    for p in tf.paragraphs:
        try:
            p.font.color.rgb = rgb
        except (AttributeError, TypeError, ValueError):
            pass
        for run in p.runs:
            try:
                run.font.color.rgb = rgb
            except (AttributeError, TypeError, ValueError):
                pass


def _apply_theme_to_chart(chart, cd, t):
    """Make chart chrome (title, axes, legend, data labels) readable on slide bg."""
    body = t["body"]
    title_rgb = t.get("title", body)
    # Tick/legend: use body (not muted) so contrast holds on dark presets
    label_rgb = body

    # Drop gallery style early — it often forces dark text on transparent plot
    try:
        chart.chart_style = None
    except (AttributeError, TypeError, ValueError):
        pass

    # Chart-wide default text (inherited where OOXML allows)
    try:
        chart.font.size = Pt(11)
        chart.font.color.rgb = label_rgb
        if t.get("font"):
            chart.font.name = t["font"]
    except (AttributeError, TypeError, ValueError):
        pass

    # Area fill → transparent over slide background
    try:
        chart.fill.background()
    except (AttributeError, TypeError, ValueError):
        pass

    # ── Title ──
    if chart.has_title:
        try:
            tf = chart.chart_title.text_frame
            _apply_rgb_to_textframe(tf, title_rgb)
            for p in tf.paragraphs:
                try:
                    p.font.size = Pt(14)
                except (AttributeError, TypeError, ValueError):
                    pass
        except (AttributeError, TypeError, ValueError):
            try:
                chart.chart_title.text_frame.paragraphs[0].font.color.rgb = title_rgb
            except (AttributeError, TypeError, ValueError):
                pass

    cat_ax = _try_category_axis(chart)
    val_ax = _try_value_axis(chart)

    x_label = cd.get("x_label", "")
    y_label = cd.get("y_label", "")

    if cat_ax is not None:
        try:
            if x_label and cat_ax.has_title:
                tf = cat_ax.axis_title.text_frame
                _apply_rgb_to_textframe(tf, label_rgb)
            cat_ax.tick_labels.font.size = Pt(10)
            cat_ax.tick_labels.font.color.rgb = label_rgb
        except (AttributeError, TypeError, ValueError):
            pass

    if val_ax is not None:
        try:
            if y_label and val_ax.has_title:
                tf = val_ax.axis_title.text_frame
                _apply_rgb_to_textframe(tf, label_rgb)
            val_ax.tick_labels.font.size = Pt(10)
            val_ax.tick_labels.font.color.rgb = label_rgb
        except (AttributeError, TypeError, ValueError):
            pass

    # Major gridlines — faint light lines on dark slides
    try:
        if val_ax is not None and val_ax.has_major_gridlines:
            val_ax.major_gridlines.format.line.color.rgb = t.get("muted", label_rgb)
    except (AttributeError, TypeError, ValueError):
        pass

    # Legend
    if chart.has_legend and chart.legend is not None:
        try:
            chart.legend.font.size = Pt(11)
            chart.legend.font.color.rgb = label_rgb
        except (AttributeError, TypeError, ValueError):
            pass

    # Data labels
    if cd.get("data_labels") is not False and chart.series:
        for s in chart.series:
            try:
                s.has_data_labels = True
                s.data_labels.font.size = Pt(9)
                s.data_labels.font.color.rgb = label_rgb
            except AttributeError:
                pass


def add_chart(slide, chart_data, left, top, width, height, t=None):
    """Add a styled OOXML chart at the given position.

    Args:
        slide: pptx slide object.
        chart_data: dict with type, categories, series, and optional
                    chart_title, x_label, y_label, data_labels.
        left, top, width, height: Inches values (EMU) for positioning.
        t: theme dict (optional). Body color applied to all chart text
           so labels remain readable on dark backgrounds.

    Returns:
        Height used (inches), or 0 if chart could not be created.
    """
    cd = chart_data
    chart_type = cd.get("type", "bar")
    is_combo = (chart_type == "combo")
    ct = XL_CHART_TYPE.COLUMN_CLUSTERED if is_combo else CHART_TYPES.get(
        chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED
    )
    categories = cd.get("categories", [])
    series_list = cd.get("series", [])

    if not categories or not series_list:
        return 0

    is_scatter = (chart_type == "scatter")
    if is_scatter:
        chart_data_obj = XyChartData()
        x_vals = [float(c) for c in categories]
        for s in series_list:
            y_vals = [float(v) for v in s.get("values", [])]
            xy_series = chart_data_obj.add_series(str(s.get("name", "")))
            for x, y in zip(x_vals, y_vals):
                xy_series.add_data_point(x, y)
    else:
        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = [str(c) for c in categories]
        for s in series_list:
            vals = [float(v) for v in s.get("values", [])]
            chart_data_obj.add_series(str(s.get("name", "")), vals)

    chart_frame = slide.shapes.add_chart(ct, left, top, width, height, chart_data_obj)
    chart = chart_frame.chart
    chart.has_legend = len(series_list) > 1 or is_combo

    # ── Combo: convert last series to line ──
    if is_combo and len(series_list) >= 2:
        # Last series → line on secondary axis
        line_series = chart.series[-1]
        line_series.chart_type = XL_CHART_TYPE.LINE_MARKERS

    chart_title_text = cd.get("chart_title", "")
    if chart_title_text:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title_text
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

    x_label = cd.get("x_label", "")
    y_label = cd.get("y_label", "")
    cat_ax = _try_category_axis(chart)
    val_ax = _try_value_axis(chart)
    if x_label and cat_ax is not None:
        cat_ax.has_title = True
        cat_ax.axis_title.text_frame.paragraphs[0].text = x_label
    if y_label and val_ax is not None:
        val_ax.has_title = True
        val_ax.axis_title.text_frame.paragraphs[0].text = y_label

    # Data labels — default ON; set "data_labels": false to suppress
    if cd.get("data_labels") is not False and chart.series:
        try:
            for s in chart.series:
                s.has_data_labels = True
                s.data_labels.font.size = Pt(9)
        except AttributeError:
            pass  # scatter (XySeries) doesn't support data_labels the same way

    # Theme: readable foreground on dark (and consistent on light) slides
    if t:
        _apply_theme_to_chart(chart, cd, t)

    return _inches(height)  # EMU → inches


def _inches(emu_val):
    return emu_val / 914400.0
