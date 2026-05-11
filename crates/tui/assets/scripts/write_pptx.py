#!/usr/bin/env python3
"""write_pptx.py — generate .pptx from JSON stdin payload.

Called by WriteOfficeTool via venv Python.
Args: --output PATH   (output .pptx path)
Stdin: JSON with "slides" array, optional "title", "subtitle", "theme".

Themes — string name (preset) or dict (custom):
  Presets: dark | light | warm | minimal
  Custom:  { "bg": "#RRGGBB", "accent": "#RRGGBB", "title": "#RRGGBB",
             "body": "#RRGGBB", "muted": "#RRGGBB", "font": "Font Name" }

Per-slide theme override: add "theme" (string or dict) to any slide data object.

Charts (per-slide, optional):
  { "chart": { "type": "bar"|"line"|"pie"|"stacked_bar"|"area"|"scatter"|"donut",
               "categories": ["Q1","Q2","Q3"],
               "series": [{ "name":"Sales","values":[10,20,30] }],
               "chart_title": "Revenue Analysis",   // optional
               "x_label": "Quarter",                // optional
               "y_label": "Revenue (10K CNY)",      // optional
               "data_labels": true } }              // optional, show values on chart

Tables (per-slide, optional):
  { "table": { "headers": ["Item","Qty","Note"],
               "rows": [["A",10,"OK"],["B",5,"WIP"]] } }

Slides can have any combination of title, bullets, chart, and table.
"""

import argparse
import json
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
except ImportError:
    print("ERROR: python-pptx not installed", file=sys.stderr)
    sys.exit(1)

# ── Hex helpers ──────────────────────────────────────────────────────────

_HEX_RE = re.compile(r"^#?([0-9a-fA-F]{6})$")


def _hex_to_rgb(s):
    m = _HEX_RE.match(str(s).strip())
    if not m:
        raise ValueError(f"invalid hex color: {s!r}")
    raw = int(m.group(1), 16)
    return RGBColor((raw >> 16) & 0xFF, (raw >> 8) & 0xFF, raw & 0xFF)


# ── Theme presets ────────────────────────────────────────────────────────

THEMES = {
    "dark": {
        "bg":     RGBColor(0x1A, 0x1A, 0x2E),
        "accent": RGBColor(0x00, 0xD4, 0xAA),
        "title":  RGBColor(0xFF, 0xFF, 0xFF),
        "body":   RGBColor(0xE0, 0xE4, 0xEC),
        "muted":  RGBColor(0x88, 0x90, 0xA0),
        "font":   "Segoe UI",
        "table_header_bg":  RGBColor(0x00, 0xD4, 0xAA),
        "table_header_fg":  RGBColor(0xFF, 0xFF, 0xFF),
        "table_row_bg":     RGBColor(0x1A, 0x1A, 0x2E),
        "table_alt_bg":     RGBColor(0x24, 0x24, 0x3A),
        "table_border":     RGBColor(0x00, 0xD4, 0xAA),
    },
    "light": {
        "bg":     RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x25, 0x63, 0xEB),
        "title":  RGBColor(0x11, 0x1F, 0x3D),
        "body":   RGBColor(0x37, 0x40, 0x51),
        "muted":  RGBColor(0x6B, 0x72, 0x80),
        "font":   "Calibri",
        "table_header_bg":  RGBColor(0x25, 0x63, 0xEB),
        "table_header_fg":  RGBColor(0xFF, 0xFF, 0xFF),
        "table_row_bg":     RGBColor(0xFF, 0xFF, 0xFF),
        "table_alt_bg":     RGBColor(0xF3, 0xF4, 0xF6),
        "table_border":     RGBColor(0x25, 0x63, 0xEB),
    },
    "warm": {
        "bg":     RGBColor(0xFF, 0xF8, 0xF0),
        "accent": RGBColor(0xE0, 0x7B, 0x3C),
        "title":  RGBColor(0x3D, 0x28, 0x1F),
        "body":   RGBColor(0x5C, 0x4A, 0x3E),
        "muted":  RGBColor(0x9C, 0x8B, 0x7E),
        "font":   "Georgia",
        "table_header_bg":  RGBColor(0xE0, 0x7B, 0x3C),
        "table_header_fg":  RGBColor(0xFF, 0xFF, 0xFF),
        "table_row_bg":     RGBColor(0xFF, 0xF8, 0xF0),
        "table_alt_bg":     RGBColor(0xFF, 0xF0, 0xE0),
        "table_border":     RGBColor(0xE0, 0x7B, 0x3C),
    },
    "minimal": {
        "bg":     RGBColor(0xFA, 0xFA, 0xFA),
        "accent": RGBColor(0x1A, 0x1A, 0x1A),
        "title":  RGBColor(0x0A, 0x0A, 0x0A),
        "body":   RGBColor(0x3D, 0x3D, 0x3D),
        "muted":  RGBColor(0x9E, 0x9E, 0x9E),
        "font":   "Helvetica",
        "table_header_bg":  RGBColor(0x1A, 0x1A, 0x1A),
        "table_header_fg":  RGBColor(0xFF, 0xFF, 0xFF),
        "table_row_bg":     RGBColor(0xFA, 0xFA, 0xFA),
        "table_alt_bg":     RGBColor(0xEE, 0xEE, 0xEE),
        "table_border":     RGBColor(0x1A, 0x1A, 0x1A),
    },
}

_THEME_KEYS = ("bg", "accent", "title", "body", "muted", "font")


def resolve_theme(raw, fallback=None):
    if raw is None:
        raw = fallback
    if raw is None:
        raw = "dark"

    if isinstance(raw, str):
        return THEMES.get(raw, THEMES["dark"])

    if isinstance(raw, dict):
        for k in _THEME_KEYS:
            if k not in raw:
                raise ValueError(f"custom theme missing key: {k}")
        base = THEMES["dark"]  # inherit table colors from dark preset
        base["bg"]     = _hex_to_rgb(str(raw["bg"]))
        base["accent"] = _hex_to_rgb(str(raw["accent"]))
        base["title"]  = _hex_to_rgb(str(raw["title"]))
        base["body"]   = _hex_to_rgb(str(raw["body"]))
        base["muted"]  = _hex_to_rgb(str(raw["muted"]))
        base["font"]   = str(raw["font"])
        return base

    raise ValueError(f"theme must be a string or dict, got {type(raw).__name__}")


# ── Layout constants ─────────────────────────────────────────────────────

SLIDE_W = 13.333
SLIDE_H = 7.5

CHART_TYPES = {
    "bar":            XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line":           XL_CHART_TYPE.LINE_MARKERS,
    "pie":            XL_CHART_TYPE.PIE,
    "stacked_bar":    XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_bar_pct": XL_CHART_TYPE.COLUMN_STACKED_100,
    "area":           XL_CHART_TYPE.AREA,
    "scatter":        XL_CHART_TYPE.XY_SCATTER,
    "donut":          XL_CHART_TYPE.DOUGHNUT,
}


# ── Drawing helpers ──────────────────────────────────────────────────────

def _add_bg(slide, clr):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = clr


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _add_para(tf, text, size, color, font, bold=False, align=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.space_after = Pt(int(size * 0.35))
    p.space_before = Pt(int(size * 0.1))
    if align is not None:
        p.alignment = align
    return p


# ── Table helper ─────────────────────────────────────────────────────────

def _add_table(slide, table_data, t):
    """Add a styled native table to the slide. Returns total height consumed."""
    td = table_data
    headers = td.get("headers", [])
    rows = td.get("rows", [])
    if not headers and not rows:
        return 0

    n_cols = len(headers) or (len(rows[0]) if rows else 1)
    n_body_rows = len(rows)
    n_total_rows = n_body_rows + (1 if headers else 0)

    max_rows = 18
    if n_total_rows > max_rows:
        n_body_rows = max_rows - (1 if headers else 0)
        rows = rows[:n_body_rows]
        n_total_rows = max_rows

    row_h = Pt(26)
    total_h = row_h * n_total_rows + Pt(2)
    top = Inches(1.85)
    left = Inches(0.6)
    width = Inches(12.133)

    tbl_shape = slide.shapes.add_table(n_total_rows, n_cols, left, top, width, total_h)
    tbl = tbl_shape.table

    # Column widths — even distribution
    col_w = int(Inches(12.133) / n_cols)
    for ci in range(n_cols):
        tbl.columns[ci].width = col_w

    # Helper: write a cell
    def _cell(ri, ci, text, color, bg_color, font_name, bold=False, size=11):
        cell = tbl.cell(ri, ci)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = PP_ALIGN.LEFT
        cell.vertical_anchor = 1  # MSO_ANCHOR.MIDDLE
        # Fill
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

    ri = 0
    # Header row
    if headers:
        for ci, h in enumerate(headers):
            _cell(ri, ci, str(h), t["table_header_fg"], t["table_header_bg"],
                  t["font"], bold=True, size=12)
        ri += 1

    # Body rows — alternating colors
    for i, row in enumerate(rows):
        bg = t["table_alt_bg"] if i % 2 else t["table_row_bg"]
        for ci in range(n_cols):
            val = str(row[ci]) if ci < len(row) else ""
            _cell(ri, ci, val, t["body"], bg, t["font"], size=11)
        ri += 1

    return total_h


# ── Chart helper ─────────────────────────────────────────────────────────

def _add_chart(slide, chart_data, top_inches, height_inches):
    """Add a styled chart at the given position. Returns actual height used."""
    cd = chart_data
    ct = CHART_TYPES.get(cd.get("type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    categories = cd.get("categories", [])
    series_list = cd.get("series", [])

    if not categories or not series_list:
        return 0

    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = [str(c) for c in categories]
    for s in series_list:
        vals = [float(v) for v in s.get("values", [])]
        chart_data_obj.add_series(str(s.get("name", "")), vals)

    left = Inches(1.2)
    top = Inches(top_inches)
    width = Inches(10.8)
    height = Inches(height_inches)

    chart_frame = slide.shapes.add_chart(ct, left, top, width, height, chart_data_obj)
    chart = chart_frame.chart
    chart.has_legend = len(series_list) > 1

    # Chart title
    chart_title_text = cd.get("chart_title", "")
    if chart_title_text:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title_text
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

    # Axis labels
    x_label = cd.get("x_label", "")
    y_label = cd.get("y_label", "")
    if x_label and chart.category_axis:
        chart.category_axis.has_title = True
        chart.category_axis.axis_title.text_frame.paragraphs[0].text = x_label
    if y_label and chart.value_axis:
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.paragraphs[0].text = y_label

    # Data labels on series
    if cd.get("data_labels") and chart.series:
        for s in chart.series:
            s.has_data_labels = True
            s.data_labels.font.size = Pt(9)

    return height_inches


# ── Slide builders ───────────────────────────────────────────────────────

def build_cover(slide, doc_title, subtitle, t):
    _add_bg(slide, t["bg"])
    _add_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W), Pt(6), t["accent"])

    top = Inches(2.2) if subtitle else Inches(2.7)
    tf = _add_textbox(slide, Inches(1.5), top, Inches(10.333), Inches(2.0))
    _add_para(tf, doc_title, 48, t["accent"], t["font"], bold=True,
              align=PP_ALIGN.CENTER, first=True)

    if subtitle:
        tf2 = _add_textbox(slide, Inches(2.0), Inches(4.5), Inches(9.333), Inches(1.0))
        _add_para(tf2, subtitle, 22, t["muted"], t["font"],
                  align=PP_ALIGN.CENTER, first=True)

    _add_rect(slide, Inches(5.0), Inches(6.8), Inches(3.333), Pt(2), t["accent"])


def build_content(slide, data, num, total, global_t):
    t = global_t
    if "theme" in data and data["theme"] is not None:
        try:
            t = resolve_theme(data["theme"], fallback=global_t)
        except ValueError:
            pass

    _add_bg(slide, t["bg"])
    _add_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W), Pt(4), t["accent"])

    title = data.get("title", "")
    if title:
        tf = _add_textbox(slide, Inches(1.0), Inches(0.55), Inches(11.333), Inches(0.9))
        _add_para(tf, title, 36, t["title"], t["font"], bold=True, first=True)
        _add_rect(slide, Inches(1.0), Inches(1.5), Inches(2.5), Pt(3), t["accent"])

    has_chart = data.get("chart") is not None
    has_table = data.get("table") is not None
    has_bullets = bool(data.get("bullets", []))

    content_top = 1.85   # inches below title accent bar
    max_avail = 5.15     # inches available for content (to leave room for bullets/notes)

    # ── Table ──
    if has_table:
        table_h_emu = _add_table(slide, data["table"], t)         # EMUs
        content_top += table_h_emu / 914400 + 0.15                # inches

    # ── Chart (below table or in remaining space) ──
    if has_chart:
        remaining = max_avail - (content_top - 1.85)
        _add_chart(slide, data["chart"], content_top, min(4.8, max(2.5, remaining - 0.1)))
        content_top += min(4.8, max(2.5, remaining - 0.1)) + 0.1

    # ── Bullets ──
    if has_bullets:
        has_content_above = has_chart or has_table
        if has_content_above:
            b_top = Inches(content_top + 0.1)
        else:
            b_top = Inches(2.0)
        b_bottom = Inches(SLIDE_H - 0.55)
        b_height = Emu(max(Inches(0.3), b_bottom - b_top))
        ft_size = 14 if not has_content_above else 12
        tf = _add_textbox(slide, Inches(1.2), b_top, Inches(10.8), b_height)
        for i, b in enumerate(data.get("bullets", [])):
            _add_para(tf, f"  {b.strip()}", ft_size, t["body"], t["font"], first=(i == 0))

    notes = data.get("notes", "")
    if notes and slide.has_notes_slide:
        slide.notes_slide.notes_text_frame.text = notes

    tf_pn = _add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4))
    _add_para(tf_pn, f"{num}/{total}", 10, t["muted"], t["font"],
              align=PP_ALIGN.RIGHT, first=True)


def build_end(slide, doc_title, t):
    _add_bg(slide, t["bg"])
    _add_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W), Pt(6), t["accent"])
    tf = _add_textbox(slide, Inches(1.5), Inches(2.4), Inches(10.333), Inches(1.5))
    _add_para(tf, "Thank You", 48, t["accent"], t["font"], bold=True,
              align=PP_ALIGN.CENTER, first=True)
    if doc_title:
        tf2 = _add_textbox(slide, Inches(2.0), Inches(4.2), Inches(9.333), Inches(0.8))
        _add_para(tf2, doc_title, 20, t["muted"], t["font"],
                  align=PP_ALIGN.CENTER, first=True)
    _add_rect(slide, Inches(5.0), Inches(6.8), Inches(3.333), Pt(2), t["accent"])


# ── Main ─────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", required=True)
    parser.add_argument("--input", default=None)
    args = parser.parse_args()

    if args.input:
        with open(args.input, "r", encoding="utf-8") as f:
            payload = json.load(f)
    else:
        payload = json.load(sys.stdin)

    try:
        t = resolve_theme(payload.get("theme"))
    except ValueError as e:
        print(f"ERROR: invalid theme — {e}", file=sys.stderr)
        sys.exit(1)

    doc_title = payload.get("title", "")
    doc_subtitle = payload.get("subtitle", "")
    slides_data = payload.get("slides", [])

    if not slides_data:
        print("ERROR: no slides data", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    total = len(slides_data) + 1
    if doc_title:
        total += 1

    if doc_title:
        build_cover(prs.slides.add_slide(prs.slide_layouts[6]),
                    doc_title, doc_subtitle, t)

    base = 1 if doc_title else 0
    for i, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        build_content(slide, sd, i + 1 + base, total, t)

    build_end(prs.slides.add_slide(prs.slide_layouts[6]), doc_title or "", t)

    prs.save(args.output)
    print("OK")


if __name__ == "__main__":
    main()
