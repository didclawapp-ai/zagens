"""pptx_engine/layout.py — slide builders and grid layout engine.

Phase 1: cover/content/end slide builders, GridLayout for blocks pipeline,
old pipeline backward compatibility, blocks pipeline orchestration.
"""

import sys

try:
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
except ImportError:
    pass

from .blocks import (
    add_bg, add_rect, add_textbox, add_para, add_table,
    parse_dim, render_block,
)
from .theme import resolve_theme

# ── Layout constants ─────────────────────────────────────────────────────

SLIDE_W = 13.333
SLIDE_H = 7.5

# Content area for blocks-pipeline pages (inches, from slide edge)
CONTENT_LEFT = 0.55
CONTENT_TOP = 1.85    # below title + accent bar
CONTENT_WIDTH = 12.233
CONTENT_HEIGHT = 5.05  # above page number area


# ── Grid layout engine ───────────────────────────────────────────────────

class GridLayout:
    """Compute column region boxes from layout config.

    Attributes:
        regions: list of (left, top, width, height) in EMU (Inches-compatible).
    """

    def __init__(self, layout_config, content_area=None):
        """Initialize from layout config.

        Args:
            layout_config: dict with 'cols' (weight list), optional 'gap', 'padding'.
            content_area: (left, top, width, height) in EMU. Defaults to standard.
        """
        if content_area is None:
            content_area = (
                Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                Inches(CONTENT_WIDTH), Inches(CONTENT_HEIGHT),
            )

        ca_left, ca_top, ca_width, ca_height = content_area
        cols = layout_config.get("cols", [1.0])
        gap = parse_dim(layout_config.get("gap", "0.2in"))
        padding = parse_dim(layout_config.get("padding", "0in"))

        # Inner area after padding
        inner_left = ca_left + padding
        inner_top = ca_top + padding
        inner_width = ca_width - 2 * padding
        inner_height = ca_height - 2 * padding

        # Total gap space
        n_cols = len(cols)
        total_gap = gap * (n_cols - 1) if n_cols > 1 else 0
        available_width = inner_width - total_gap

        # Weight sum
        total_weight = sum(cols)
        if total_weight <= 0:
            total_weight = 1.0

        # Compute regions
        self.regions = []
        cursor = inner_left
        for i, w in enumerate(cols):
            col_w = int(available_width * w / total_weight)
            region = (cursor, inner_top, col_w, inner_height)
            self.regions.append(region)
            cursor += col_w + gap


# ── Slide builders ───────────────────────────────────────────────────────

def build_cover(slide, doc_title, subtitle, t):
    add_bg(slide, t["bg"])
    add_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W), Pt(6), t["accent"])

    top = Inches(2.2) if subtitle else Inches(2.7)
    tf = add_textbox(slide, Inches(1.5), top, Inches(10.333), Inches(2.0))
    add_para(tf, doc_title, 48, t["accent"], t["font"], bold=True,
             align=PP_ALIGN.CENTER, first=True)

    if subtitle:
        tf2 = add_textbox(slide, Inches(2.0), Inches(4.5), Inches(9.333), Inches(1.0))
        add_para(tf2, subtitle, 22, t["muted"], t["font"],
                 align=PP_ALIGN.CENTER, first=True)

    add_rect(slide, Inches(5.0), Inches(6.8), Inches(3.333), Pt(2), t["accent"])


def build_content(slide, data, num, total, global_t):
    """Build a content slide using the old pipeline (chart→table→bullets).

    This is the backward-compatible path.  When 'blocks' is present
    and non-empty, build_content_blocks() is used instead (§5.0).
    """
    from .charts import add_chart  # only needed for old pipeline

    t = global_t
    if "theme" in data and data["theme"] is not None:
        try:
            t = resolve_theme(data["theme"], fallback=global_t)
        except ValueError:
            pass

    add_bg(slide, t["bg"])
    add_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W), Pt(4), t["accent"])

    title = data.get("title", "")
    if title:
        tf = add_textbox(slide, Inches(1.0), Inches(0.55), Inches(11.333), Inches(0.9))
        add_para(tf, title, 36, t["title"], t["font"], bold=True, first=True)
        add_rect(slide, Inches(1.0), Inches(1.5), Inches(2.5), Pt(3), t["accent"])

    has_chart = data.get("chart") is not None
    has_table = data.get("table") is not None
    has_bullets = bool(data.get("bullets", []))

    content_top = 1.85   # inches below title accent bar
    max_avail = 5.15

    # Table — Note: add_table returns EMU; /914400 converts to inches (legacy;
    # new blocks pipeline uses grid layout which is pure EMU).
    if has_table:
        table_h_emu = add_table(slide, data["table"], t)
        content_top += table_h_emu / 914400 + 0.15

    # Chart
    if has_chart:
        remaining = max_avail - (content_top - 1.85)
        chart_h = min(4.8, max(2.5, remaining - 0.1))
        add_chart(slide, data["chart"],
                  Inches(1.2), Inches(content_top),
                  Inches(10.8), Inches(chart_h), t=t)
        content_top += chart_h + 0.1

    # Bullets
    if has_bullets:
        has_content_above = has_chart or has_table
        if has_content_above:
            b_top = Inches(content_top + 0.1)
        else:
            b_top = Inches(2.0)
        b_bottom = Inches(SLIDE_H - 0.55)
        b_height = Emu(max(Inches(0.3), b_bottom - b_top))
        ft_size = 14 if not has_content_above else 12
        tf = add_textbox(slide, Inches(1.2), b_top, Inches(10.8), b_height)
        for i, b in enumerate(data.get("bullets", [])):
            add_para(tf, f"  {b.strip()}", ft_size, t["body"], t["font"], first=(i == 0))

    _add_notes_and_page_num(slide, data, num, total, t)


def build_end(slide, doc_title, t):
    add_bg(slide, t["bg"])
    add_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W), Pt(6), t["accent"])
    tf = add_textbox(slide, Inches(1.5), Inches(2.4), Inches(10.333), Inches(1.5))
    add_para(tf, "Thank You", 48, t["accent"], t["font"], bold=True,
             align=PP_ALIGN.CENTER, first=True)
    if doc_title:
        tf2 = add_textbox(slide, Inches(2.0), Inches(4.2), Inches(9.333), Inches(0.8))
        add_para(tf2, doc_title, 20, t["muted"], t["font"],
                 align=PP_ALIGN.CENTER, first=True)
    add_rect(slide, Inches(5.0), Inches(6.8), Inches(3.333), Pt(2), t["accent"])


# ── Blocks pipeline ──────────────────────────────────────────────────────

def build_content_blocks(slide, data, num, total, global_t):
    """Build a content slide using the blocks pipeline with grid layout.

    Falls back to old pipeline on any error (§5.0).
    """
    t = global_t
    if "theme" in data and data["theme"] is not None:
        try:
            t = resolve_theme(data["theme"], fallback=global_t)
        except ValueError:
            pass

    try:
        _build_blocks_inner(slide, data, num, total, t)
    except Exception as e:
        print(
            f"WARNING: slide {num} blocks pipeline failed ({e}); "
            f"falling back to old pipeline",
            file=sys.stderr,
        )
        build_content(slide, data, num, total, global_t)


def _build_blocks_inner(slide, data, num, total, t):
    """Inner blocks pipeline — may raise on schema errors."""
    # Background + title bar chrome
    add_bg(slide, t["bg"])
    add_rect(slide, Inches(0), Inches(0), Inches(SLIDE_W), Pt(4), t["accent"])

    title = data.get("title", "")
    if title:
        tf = add_textbox(slide, Inches(1.0), Inches(0.55), Inches(11.333), Inches(0.9))
        add_para(tf, title, 36, t["title"], t["font"], bold=True, first=True)
        add_rect(slide, Inches(1.0), Inches(1.5), Inches(2.5), Pt(3), t["accent"])

    blocks = data["blocks"]
    layout_config = data.get("layout", {})

    grid = GridLayout(layout_config)
    regions = grid.regions

    # Dispatch blocks to grid regions
    for i, block in enumerate(blocks):
        if i >= len(regions):
            print(
                f"WARNING: slide {num} has more blocks ({len(blocks)}) "
                f"than regions ({len(regions)}); truncating",
                file=sys.stderr,
            )
            break
        left, top, width, height = regions[i]
        render_block(slide, block, left, top, width, height, t)

    _add_notes_and_page_num(slide, data, num, total, t)


def _footer_text_color(t):
    """Readable page number on dark slides (muted alone is often too dim)."""
    bg = t["bg"]
    lum = (bg[0] * 299 + bg[1] * 587 + bg[2] * 114) / 1000.0
    if lum < 140:
        return t["body"]
    return t.get("muted", t["body"])


# ── Shared helpers ───────────────────────────────────────────────────────

def _add_notes_and_page_num(slide, data, num, total, t):
    """Add speaker notes and page number footer."""
    notes = data.get("notes", "")
    if notes and slide.has_notes_slide:
        slide.notes_slide.notes_text_frame.text = notes

    tf_pn = add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4))
    add_para(tf_pn, f"{num}/{total}", 10, _footer_text_color(t), t["font"],
             align=PP_ALIGN.RIGHT, first=True)
